import streamlit as st
from io import BytesIO
from pptx import Presentation

# Function to extract text from PPTX and remove spaces left by images
def extract_text_from_pptx(file):
    """Extract text from a PowerPoint file and remove empty spaces where images are placed."""
    presentation = Presentation(file)
    text = []
    
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Strip any extra spaces and add to slide_text
                slide_text.append(shape.text.strip())
            # Skip image shapes (no text content)
            elif shape.shape_type == 13:  # Type 13 is for pictures (images)
                continue
        
        # If there is any text in the slide, add it to the main text, else ignore the slide
        if slide_text:
            text.append("\n".join(slide_text))
    
    return "\n".join(text)

# Set page layout
st.set_page_config(page_title="PPT Text Extractor", layout="centered")

# Style and layout for the landing page
landing_page_style = """
    <style>
        .title {
            text-align: center;
            color: #2F8F4F;
            font-size: 36px;
            font-family: 'Arial', sans-serif;
        }
        .subtitle {
            text-align: center;
            color: #777;
            font-size: 18px;
            font-family: 'Arial', sans-serif;
        }
        .description {
            font-size: 18px;
            color: #444;
            text-align: center;
            font-family: 'Arial', sans-serif;
        }
        .button-container {
            display: flex;
            justify-content: center;
            margin-top: 20px;
        }
        .footer {
            text-align: center;
            font-size: 14px;
            color: #888;
            margin-top: 40px;
            font-family: 'Arial', sans-serif;
        }
        .footer a {
            color: #333;
            text-decoration: none;
            font-size: 16px;
        }
        .footer a:hover {
            color: #4CAF50;
        }
    </style>
"""

# Apply custom styles to the landing page
st.markdown(landing_page_style, unsafe_allow_html=True)

# Title and Description
st.markdown('<div class="title">PowerPoint Text Extractor</div>', unsafe_allow_html=True)
st.markdown('<div class="subtitle">Upload a PowerPoint (.pptx) file to extract its text content.</div>', unsafe_allow_html=True)

# File uploader widget
uploaded_file = st.file_uploader("Choose a PPT file", type=["pptx"])

# If a file is uploaded, process it
if uploaded_file is not None:
    # Read the file
    file_bytes = uploaded_file.read()
    
    # Extract text from the uploaded PPT
    text = extract_text_from_pptx(BytesIO(file_bytes))
    
    # Display the extracted text
    st.subheader("Extracted Text")

    # Display the extracted text in a text area
    st.text_area("Text Content", text, height=300, key="extracted_text", disabled=True)
    
    # Display the download button in the middle of the screen
    st.markdown('<div class="button-container">', unsafe_allow_html=True)
    st.download_button(
        label="Download Extracted Text",
        data=text,
        file_name="extracted_text.txt",
        mime="text/plain",
        use_container_width=True
    )
    st.markdown('</div>', unsafe_allow_html=True)

    # Show message for successful extraction
    st.markdown("""
        <div style="text-align:center; margin-top:20px;">
            <p style="color: #4CAF50; font-size: 18px;">Text successfully extracted!</p>
        </div>
    """, unsafe_allow_html=True)

# Footer with name, email, and WhatsApp in one line
st.markdown('<div class="footer">', unsafe_allow_html=True)
st.markdown('<span>Developed by: Muhammad Haris | Email me: imharis.dev@gmail.com | Whatsapp me: 0318-8466565 </span>', unsafe_allow_html=True)
st.markdown('</div>', unsafe_allow_html=True)
